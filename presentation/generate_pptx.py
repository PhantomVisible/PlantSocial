from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Color Palette ──────────────────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0D, 0x1B, 0x0E)
C_BG_CARD   = RGBColor(0x16, 0x25, 0x18)
C_GREEN     = RGBColor(0x2E, 0x7D, 0x32)
C_GREEN_LT  = RGBColor(0x4C, 0xAF, 0x50)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED     = RGBColor(0x8A, 0xAB, 0x8C)
C_ACCENT    = RGBColor(0x81, 0xC7, 0x84)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ─── Helper functions ────────────────────────────────────────────────────────

def add_bg(slide, color):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txb(slide, text, l, t, w, h,
        size=Pt(14), bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def green_tag(slide, text, l, t, w=2.5):
    """Small green tag label."""
    add_rect(slide, l, t, w, 0.28, fill_color=RGBColor(0x1A,0x3D,0x1C), line_color=C_GREEN_LT, line_width=Pt(0.5))
    txb(slide, text.upper(), l+0.08, t+0.02, w-0.1, 0.24,
        size=Pt(7), bold=True, color=C_GREEN_LT, align=PP_ALIGN.CENTER)

def slide_title(slide, text, l, t, w=11.5, size=Pt(36)):
    """Main title with accent on last word."""
    txb(slide, text, l, t, w, 0.7, size=size, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def subtitle_txt(slide, text, l, t, w=9):
    txb(slide, text, l, t, w, 0.4, size=Pt(13), color=C_MUTED, align=PP_ALIGN.CENTER)

def card(slide, l, t, w, h, title, desc, icon=""):
    """Draw a card with a top green border."""
    add_rect(slide, l,     t,       w,    0.04,  fill_color=C_GREEN_LT)
    add_rect(slide, l, t+0.04, w, h-0.04, fill_color=C_BG_CARD, line_color=RGBColor(0x2A,0x3E,0x2C), line_width=Pt(0.5))
    if icon:
        txb(slide, icon, l+0.12, t+0.12, 0.5, 0.35, size=Pt(18))
    txb(slide, title, l+0.65 if icon else l+0.12, t+0.12, w-0.75, 0.25, size=Pt(10), bold=True, color=C_WHITE)
    txb(slide, desc,  l+0.12, t+0.44, w-0.24, h-0.55, size=Pt(8.5), color=C_MUTED)

def flow_arrow(slide, l, t):
    txb(slide, "→", l, t, 0.3, 0.35, size=Pt(18), color=C_GREEN_LT, align=PP_ALIGN.CENTER)

def flow_box(slide, l, t, label, sub=""):
    add_rect(slide, l, t, 1.55, 0.65, fill_color=C_BG_CARD, line_color=RGBColor(0x2A,0x3E,0x2C), line_width=Pt(0.5))
    if sub:
        txb(slide, sub.upper(), l+0.08, t+0.06, 1.4, 0.2, size=Pt(6.5), bold=True, color=C_GREEN_LT)
    txb(slide, label, l+0.08, t+0.27, 1.4, 0.32, size=Pt(9), bold=True, color=C_WHITE)

def quote_box(slide, text, l, t, w, h):
    add_rect(slide, l, t, 0.04, h, fill_color=C_GREEN_LT)
    add_rect(slide, l+0.04, t, w-0.04, h, fill_color=RGBColor(0x1A,0x2E,0x1B), line_color=RGBColor(0x2A,0x45,0x2C), line_width=Pt(0.5))
    txb(slide, text, l+0.18, t+0.1, w-0.3, h-0.2, size=Pt(10), color=C_ACCENT, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: COVER
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)
add_rect(s, 0.5, 2.8, 3.5, 3.5, fill_color=RGBColor(0x18,0x35,0x1A))  # decorative blob
green_tag(s, "Projet Final — 2026", 5.2, 0.6, 3.0)
txb(s, "PlantSocial", 1.0, 1.1, 11, 1.4, size=Pt(72), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txb(s, "Le premier réseau social et marketplace dédié aux passionnés de plantes", 2.0, 2.7, 9, 0.5, size=Pt(14), color=C_MUTED, align=PP_ALIGN.CENTER)
members = ["👤 Membre 1 — Introduction", "👤 Membre 2 — Chat & Marketplace", "👤 Membre 3 — UI/UX & Démo", "👤 Membre 4 — Business & IA"]
for i, m in enumerate(members):
    ml = 1.0 + i*3.1
    add_rect(s, ml, 3.5, 2.9, 0.38, fill_color=C_BG_CARD, line_color=RGBColor(0x2A,0x3E,0x2C), line_width=Pt(0.5))
    txb(s, m, ml+0.1, 3.55, 2.7, 0.3, size=Pt(9), color=C_WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: AGENDA
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)
green_tag(s, "Sommaire", 5.5, 0.3, 2.2)
slide_title(s, "Plan de la présentation", 0.9, 0.7)
agenda_items = [
    ("🔍","01 — Introduction","Contexte, problématique et opportunité de marché"),
    ("🎯","02 — Objectifs","Ce que nous voulions accomplir"),
    ("⚙️","03 — Architecture","Stack technique et choix d'architecture"),
    ("🌿","04 — Fonctionnalités","Chat, Notifications, Marketplace"),
    ("🪴","05 — My Garden & IA","Identification de plante via Pl@ntNet"),
    ("🩺","06 — Doctor Plant","Diagnostic santé de la plante par IA"),
    ("📰","07 — Section News","Actualités plantes intégrées dans l'app"),
    ("💡","08 — Démo Live","Démonstration complète de l'application"),
    ("📈","09 — Business & Avenir","Modèle économique et perspectives futures"),
]
cols, rows = 3, 3
cw, ch = 3.9, 1.35
margin_l, margin_t = 0.7, 1.55
for idx, (icon, title, desc) in enumerate(agenda_items):
    col = idx % cols
    row = idx // cols
    cl = margin_l + col*(cw+0.15)
    ct = margin_t + row*(ch+0.12)
    card(s, cl, ct, cw, ch, title, desc, icon)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: INTRODUCTION
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_CARD)
green_tag(s, "01 — Introduction", 4.7, 0.3, 3.6)
slide_title(s, "Le problème à résoudre", 0.9, 0.7, size=Pt(34))
problems = [
    ("📱","Communauté dispersée","Les passionnés s'éparpillent entre Facebook, Instagram, Forums sans cohérence ni plateforme dédiée"),
    ("🛒","Achat/Vente non spécialisé","Leboncoin et Facebook Marketplace ne sont pas adaptés à la vente de plantes rares & collections"),
    ("🌱","Absence de confiance","Impossible de vérifier l'expertise d'un vendeur — risque fort pour les plantes rares à haute valeur"),
]
for i, (icon, title, desc) in enumerate(problems):
    t = 1.5 + i*1.45
    add_rect(s, 0.7, t, 12, 1.25, fill_color=C_BG_DARK, line_color=RGBColor(0x2A,0x3E,0x2C), line_width=Pt(0.5))
    txb(s, icon, 0.9, t+0.3, 0.6, 0.6, size=Pt(22))
    txb(s, title, 1.65, t+0.12, 4, 0.32, size=Pt(11), bold=True, color=C_WHITE)
    txb(s, desc, 1.65, t+0.47, 10.5, 0.6, size=Pt(9.5), color=C_MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: OBJECTIVES
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)
green_tag(s, "02 — Objectifs", 5.2, 0.3, 2.9)
slide_title(s, "Nos objectifs", 0.9, 0.7)
objs = [
    ("🤝","Créer une communauté","Un réseau social dédié pour partager, apprendre et s'inspirer autour des plantes"),
    ("🛍️","Marketplace de confiance","Un espace d'achat/vente sécurisé, rattaché aux profils sociaux des vendeurs"),
    ("⚡","Instantanéité totale","Chat et notifications en temps réel pour des interactions naturelles et engageantes"),
    ("🎨","UX Premium","Interface élégante, responsive, avec mode sombre/clair pour tous les appareils"),
]
cw, ch = 6.0, 2.2
for i, (icon, title, desc) in enumerate(objs):
    cl = 0.7 + (i%2)*(cw+0.3)
    ct = 1.5 + (i//2)*(ch+0.2)
    card(s, cl, ct, cw, ch, title, desc, icon)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5: ARCHITECTURE OVERVIEW
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_CARD)
green_tag(s, "03 — Architecture Technique", 4.0, 0.3, 5.0)
slide_title(s, "Stack & Architecture", 0.9, 0.7)
subtitle_txt(s, "Architecture Monolithique Modulaire — robuste, maintenable, scalable", 2.2, 1.35)
pills = [("☕","Spring Boot"),("🅰","Angular"),("🐘","PostgreSQL"),("🔐","JWT Auth"),("🔌","WebSockets"),("🐳","Docker")]
pw = 1.9
for i, (icon, name) in enumerate(pills):
    pl = 0.85 + i*(pw+0.18)
    add_rect(s, pl, 1.9, pw, 0.9, fill_color=C_BG_DARK, line_color=RGBColor(0x2A,0x3E,0x2C), line_width=Pt(0.5))
    txb(s, icon, pl+0.2, 1.96, 0.5, 0.4, size=Pt(18))
    txb(s, name, pl+0.7, 2.06, pw-0.8, 0.3, size=Pt(9), bold=True, color=C_MUTED)
# flow diagram
boxes = [("Frontend","Angular SPA"), ("API REST + WS","Spring Boot"), ("Base de données","PostgreSQL")]
for i, (sub, lbl) in enumerate(boxes):
    fl = 1.3 + i*3.65
    flow_box(s, fl, 3.15, lbl, sub)
    if i < 2:
        flow_arrow(s, fl+1.57, 3.28)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6: BACKEND DETAIL
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)
green_tag(s, "03 — Architecture : Backend", 4.2, 0.3, 4.6)
slide_title(s, "Backend — Spring Boot", 0.9, 0.7, size=Pt(32))
subtitle_txt(s, "Architecture en couches claire et totalement modulaire", 2.5, 1.32)
backend_items = [
    ("🔐","Sécurité — JWT","Authentification stateless via JSON Web Tokens. Chaque requête est vérifiée par Spring Security."),
    ("🌐","REST Controllers","Endpoints par domaine : /auth, /users, /marketplace, /chat, /feed, /plants"),
    ("⚙️","Service Layer","Logique métier isolée dans des Services Spring. Chaque fonctionnalité a son propre service."),
    ("🗄️","JPA / Repository","Accès à PostgreSQL via Spring Data JPA. Requêtes @Query pour les cas complexes."),
    ("🔌","WebSocket — STOMP","Protocole STOMP gère les abonnements et la diffusion de messages en temps réel."),
    ("🌍","APIs Externes","Intégration Pl@ntNet API + NewsAPI via RestTemplate / WebClient."),
]
cw, ch = 3.9, 1.55
ml, mt = 0.7, 1.7
for idx, (icon, title, desc) in enumerate(backend_items):
    cl = ml + (idx%3)*(cw+0.15)
    ct = mt + (idx//3)*(ch+0.12)
    card(s, cl, ct, cw, ch, title, desc, icon)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7: FRONTEND DETAIL
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_CARD)
green_tag(s, "03 — Architecture : Frontend", 4.2, 0.3, 4.6)
slide_title(s, "Frontend — Angular", 0.9, 0.7, size=Pt(32))
subtitle_txt(s, "Application Angular organisée en Feature Modules, totalement lazy-loaded", 2.0, 1.32)
frontend_items = [
    ("📦","Feature Modules","Chaque fonctionnalité (Auth, Feed, Marketplace, Chat, Profile) est un module Angular isolé."),
    ("🔄","Services & Observables","APIs REST via HttpClient. Données async gérées avec RxJS Observables et async pipe."),
    ("🛡️","Guards & Interceptors","Route Guards protègent les pages privées. HttpInterceptor injecte le JWT automatiquement."),
    ("🎨","SCSS & Theme System","Variables CSS dynamiques (--text-color, --surface-card) pour mode clair/sombre."),
    ("⚡","WebSocket Client","Connexion STOMP/SockJS côté client. Gestion de la connexion et abonnements aux topics."),
    ("🧩","PrimeNG UI Library","Composants UI de PrimeNG personnalisés avec SCSS pour une apparence professionnelle."),
]
cw, ch = 3.9, 1.55
ml, mt = 0.7, 1.7
for idx, (icon, title, desc) in enumerate(frontend_items):
    cl = ml + (idx%3)*(cw+0.15)
    ct = mt + (idx//3)*(ch+0.12)
    card(s, cl, ct, cw, ch, title, desc, icon)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8: CHAT & NOTIFICATIONS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)
green_tag(s, "04 — Ma Contribution : Temps Réel", 3.8, 0.3, 5.5)
slide_title(s, "Chat & Notifications en temps réel", 0.9, 0.7, size=Pt(30))
card(s, 0.7, 1.55, 5.8, 2.1, "Real-Time Chat (WebSockets)", "Tunnel de communication bidirectionnel et permanent. Aucune latence, aucun rechargement. Messages livrés en moins d'une milliseconde.", "💬")
card(s, 0.7, 3.8, 5.8, 2.1, "Notifications Push", "Alertes instantanées pour les messages reçus, commentaires, likes et interactions sur les annonces.", "🔔")
quote_box(s, "🔌 HTTP classique : le client demande → le serveur répond.\n\n⚡ WebSocket : connexion toujours ouverte. Le serveur pousse les mises à jour instantanément. C'est ce qui donne vie au chat en temps réel.", 6.85, 1.55, 5.8, 4.35)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9: MARKETPLACE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_CARD)
green_tag(s, "04 — Ma Contribution : Marketplace", 3.8, 0.3, 5.5)
slide_title(s, "La Marketplace intégrée", 0.9, 0.7, size=Pt(32))
mkt_items = [
    ("📋","Création d'annonces avancées","Galerie photos multiple, détection automatique du titre/prix via URL, gestion des statuts (Actif, Expiré, Brouillon)"),
    ("🛡️","Confiance par le profil social","Chaque annonce est liée au profil public du vendeur — ses posts, ses plantes, sa réputation dans la communauté"),
    ("💳","Simulation de paiement réaliste","Animation de paiement séquentielle (insertion de carte → scan → succès) pour une expérience premium"),
]
for i, (icon, title, desc) in enumerate(mkt_items):
    t = 1.5 + i*1.6
    add_rect(s, 0.7, t, 12, 1.4, fill_color=C_BG_DARK, line_color=RGBColor(0x2A,0x3E,0x2C), line_width=Pt(0.5))
    txb(s, icon, 0.9, t+0.35, 0.6, 0.6, size=Pt(22))
    txb(s, title, 1.65, t+0.12, 4, 0.32, size=Pt(11), bold=True, color=C_WHITE)
    txb(s, desc, 1.65, t+0.48, 10.5, 0.6, size=Pt(9.5), color=C_MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10: FEATURES OVERVIEW
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)
green_tag(s, "04 — Fonctionnalités Globales", 4.1, 0.3, 5.0)
slide_title(s, "Ce que PlantSocial offre", 0.9, 0.7, size=Pt(32))
feat_items = [
    ("👤","Profils Utilisateurs","Collection de plantes, posts, historique et réputation publique"),
    ("📰","Fil d'Actualité","Posts, commentaires, likes et partages en temps réel"),
    ("🌙","Mode Sombre / Clair","Interface entièrement adaptée aux deux modes pour le confort visuel"),
    ("🔒","Authentification sécurisée","JWT, validation stricte des mots de passe, autorisation par rôle"),
    ("📱","100% Responsive","Design optimisé mobile, tablette et desktop"),
    ("🖼️","Galerie Photos Shopify-style","Thumbnails latéraux, image principale grande, navigation fluide"),
]
cw, ch = 3.9, 1.55
ml, mt = 0.7, 1.7
for idx, (icon, title, desc) in enumerate(feat_items):
    cl = ml + (idx%3)*(cw+0.15)
    ct = mt + (idx//3)*(ch+0.12)
    card(s, cl, ct, cw, ch, title, desc, icon)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11: MY GARDEN + Pl@ntNet
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_CARD)
green_tag(s, "05 — Intégration IA : My Garden", 3.8, 0.3, 5.5)
slide_title(s, "Mon Jardin — Pl@ntNet API", 0.9, 0.7, size=Pt(32))
subtitle_txt(s, "Identification automatique de plantes grâce à l'Intelligence Artificielle", 2.0, 1.32)
flow_items = [("📸\nUpload Photo","Étape 1"), ("🌐\nPl@ntNet API","Étape 2"), ("🌿\nNom Identifié","Étape 3"), ("💾\nSauvegarde BDD","Étape 4")]
fl = 0.9
for i, (lbl, sub) in enumerate(flow_items):
    flow_box(s, fl, 2.0, lbl, sub)
    fl += 1.57
    if i < 3:
        flow_arrow(s, fl, 2.15)
        fl += 0.35
card(s, 0.7, 3.1, 6.0, 2.1, "Collection Personnelle", "L'utilisateur peut photographier n'importe quelle plante. L'IA retourne son nom scientifique et commun instantanément, puis l'ajoute à son jardin virtuel.", "🪴")
card(s, 6.9, 3.1, 6.0, 2.1, "API Pl@ntNet Intégrée", "Pl@ntNet est une IA botanique capable d'identifier +300 000 espèces de plantes. Connectée directement à notre backend Spring Boot.", "🧠")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12: DOCTOR PLANT
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)
green_tag(s, "06 — Intégration IA : Doctor Plant", 3.7, 0.3, 5.7)
slide_title(s, "Doctor Plant — Diagnostic par IA", 0.9, 0.7, size=Pt(30))
subtitle_txt(s, "Votre plante est malade ? L'IA la diagnostique avant même un expert !", 2.0, 1.32)
card(s, 0.7, 1.85, 5.8, 1.8, "Upload d'une Photo", "L'utilisateur prend en photo les feuilles, la tige, ou les racines de sa plante souffrante et l'uploade dans l'application.", "📸")
card(s, 0.7, 3.8, 5.8, 1.8, "Analyse par l'IA", "L'IA analyse l'image et identifie les symptômes visuels : jaunissement, pourriture, parasites, manque d'eau ou de lumière.", "🤖")
card(s, 6.85, 1.85, 5.8, 1.8, "Rapport de Santé", "L'application retourne un rapport détaillé : diagnostic du problème, niveau de gravité, et recommandations de traitement.", "📋")
quote_box(s, "💡 Valeur ajoutée clé : Un diagnostic instantané qui transforme chaque utilisateur en expert botanique — une fonctionnalité impossible à trouver sur une plateforme généraliste.", 6.85, 3.8, 5.8, 1.8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13: NEWS SECTION
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_CARD)
green_tag(s, "07 — Section News", 5.1, 0.3, 3.1)
slide_title(s, "Actualités Plantes intégrées", 0.9, 0.7, size=Pt(32))
subtitle_txt(s, "Restez informés des dernières tendances botaniques sans quitter PlantSocial", 2.0, 1.32)
news_items = [
    ("📡","News API intégrée","Nous consommons une API d'actualités externe pour agréger les dernières nouvelles du monde botanique"),
    ("🌿","Contenu Ciblé","Articles filtrés : conseils d'entretien, nouvelles espèces, tendances jardinage, événements horticoles"),
    ("🔗","Engagement Renforcé","Garder l'utilisateur actif même sans interactions sociales — toujours quelque chose de nouveau à découvrir"),
]
cw, ch = 3.9, 2.0
ml = 0.7
for idx, (icon, title, desc) in enumerate(news_items):
    cl = ml + idx*(cw+0.22)
    card(s, cl, 1.85, cw, ch, title, desc, icon)
quote_box(s, "📰 La section News transforme PlantSocial en hub d'information botanique — pas seulement un réseau social mais une source de connaissance vivante sur les plantes.", 0.7, 4.05, 12.0, 1.0)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14: BUSINESS MODEL
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)
green_tag(s, "09 — Business & Avenir", 4.7, 0.3, 3.8)
slide_title(s, "Modèle Économique & Perspectives", 0.9, 0.7, size=Pt(30))
biz_items = [
    ("🆓","Réseau Social Gratuit","Acquisition massive d'utilisateurs grâce à la gratuité totale de la plateforme sociale"),
    ("💰","Listing Fee","Frais de mise en ligne proportionnels à la durée d'exposition sur la Marketplace"),
    ("🏪","Abonnement Pro","Compte Pro pour pépinières et marchands avec outils avancés et visibilité boostée"),
    ("🪴","My Garden Premium","Identification illimitée de plantes & accès complet Pl@ntNet en abonnement mensuel"),
    ("🩺","Doctor Plant Pro","Diagnostics illimités + suivi médical de la collection pour les utilisateurs les plus investis"),
    ("📣","Publicité Ciblée","Audience 100% qualifiée = publicités ultra-pertinentes pour engrais, pots, lampes UV"),
]
cw, ch = 3.9, 1.55
ml, mt = 0.7, 1.65
for idx, (icon, title, desc) in enumerate(biz_items):
    cl = ml + (idx%3)*(cw+0.15)
    ct = mt + (idx//3)*(ch+0.12)
    card(s, cl, ct, cw, ch, title, desc, icon)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15: DEMO LIVE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_GREEN)
add_rect(s, 0, 0, 13.33, 7.5, fill_color=RGBColor(0x1B,0x5E,0x20))
green_tag(s, "Démonstration Live", 4.9, 0.4, 3.5)
slide_title(s, "Démonstration Live", 0.9, 0.9, size=Pt(40))
subtitle_txt(s, "Parcours utilisateur complet — de l'inscription au diagnostic IA", 2.0, 1.65)
demo_steps = [
    ("📋 Inscription","Étape 1"),
    ("📰 Feed & Profil","Étape 2"),
    ("💬 Chat Live","Étape 3"),
    ("🛒 Marketplace","Étape 4"),
    ("🪴 My Garden","Étape 5"),
    ("🩺 Doctor Plant","Étape 6"),
]
fl = 0.5
for i, (lbl, sub) in enumerate(demo_steps):
    add_rect(s, fl, 2.45, 1.85, 0.8, fill_color=RGBColor(0x00,0x00,0x00))
    shape = s.shapes[-1]
    shape.fill.fore_color.rgb = RGBColor(0x00,0x00,0x00)
    shape.fill.transparency = 0.7
    add_rect(s, fl, 2.45, 1.85, 0.8, fill_color=None, line_color=RGBColor(0xC8,0xE6,0xC9), line_width=Pt(0.5))
    txb(s, sub.upper(), fl+0.08, 2.5, 1.7, 0.18, size=Pt(6.5), bold=True, color=C_ACCENT)
    txb(s, lbl, fl+0.08, 2.72, 1.7, 0.38, size=Pt(9), bold=True, color=C_WHITE)
    fl += 1.87
    if i < 5:
        txb(s, "→", fl, 2.65, 0.25, 0.35, size=Pt(14), color=C_WHITE, align=PP_ALIGN.CENTER)
        fl += 0.28

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16: CONCLUSION
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, C_BG_DARK)
add_rect(s, 0.5, 3.0, 3.0, 3.0, fill_color=RGBColor(0x18,0x35,0x1A))
green_tag(s, "Conclusion", 5.5, 0.3, 2.2)
slide_title(s, "PlantSocial — Une vision réelle", 0.9, 0.7, size=Pt(32))
stats = [("4","Développeurs"),("6+","Modules"),("⚡","Temps Réel"),("🌿","Passion")]
sw = 2.8
for i, (num, label) in enumerate(stats):
    sl = 0.9 + i*(sw+0.35)
    add_rect(s, sl, 1.75, sw, 1.5, fill_color=C_BG_CARD, line_color=RGBColor(0x2A,0x3E,0x2C), line_width=Pt(0.5))
    txb(s, num, sl, 1.88, sw, 0.7, size=Pt(32), bold=True, color=C_GREEN_LT, align=PP_ALIGN.CENTER)
    txb(s, label.upper(), sl, 2.62, sw, 0.3, size=Pt(8), color=C_MUTED, align=PP_ALIGN.CENTER)
quote_box(s, "🙏 Nous vous remercions chaleureusement pour votre attention.\nPlantSocial est né d'un vrai problème, construit sur de vraies technologies, et vise un vrai marché en croissance.\n\nNous sommes à présent disponibles pour vos questions.", 0.9, 3.55, 11.6, 1.6)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out = "PlantSocial_Presentation.pptx"
prs.save(out)
print(f"✅ Saved: {out}")
